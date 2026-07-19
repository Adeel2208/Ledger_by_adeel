"""Pitch-deck connector (FR-4): PDF/PPTX/DOCX -> per-slide text -> structured claims.

Two stages:
  1. EXTRACT — pull raw text per slide/page (source tagged "deck:slide-N" so every
     claim traces back to an exact slide for evidence / agentic traceability, I1).
  2. COGNIFY — an LLM turns each slide's text into structured `ExtractedClaim`s.
     Deck claims are `claimed` confidence (founder-asserted) until corroborated.

Never fabricate a missing metric — the extractor is instructed to leave unknowns null.
"""
from __future__ import annotations

import os
from datetime import datetime, timezone
from typing import Any

from app.llm.client import get_llm
from app.memory.ingestion.base import BaseConnector, IngestBundle, RawSignal
from app.schemas.common import Confidence
from app.schemas.founder import ExtractedProfile

_EXTRACT_SYSTEM = (
    "You extract structured facts from a startup pitch deck for an investor's data "
    "system. Only record what is explicitly stated. If a field is not present, leave "
    "it null — never guess or infer numbers. Attribute each claim to its slide."
)


class DeckConnector(BaseConnector):
    name = "deck"

    def fetch(self, *, file_path: str, **_: Any) -> IngestBundle:
        slides = self._extract_text(file_path)
        if not slides:
            return IngestBundle(signals=[])

        full_text = "\n\n".join(f"[slide {i + 1}]\n{t}" for i, t in enumerate(slides))
        profile = self._cognify(full_text)

        now = datetime.now(timezone.utc)
        signals: list[RawSignal] = []
        for claim in profile.claims:
            slide_ref = claim.location or "deck"
            signals.append(
                RawSignal(
                    source=f"deck:{slide_ref}",
                    record_type=claim.record_type,
                    payload={"text": claim.text, "value": claim.value, "slide": slide_ref},
                    timestamp=now,
                    confidence=claim.confidence or Confidence.CLAIMED,
                    external_url=None,
                    text=claim.text,
                )
            )

        return IngestBundle(
            signals=signals,
            identity={
                "name": profile.founder_name,
                "email": profile.founder_email,
                "github_handle": profile.github_handle,
                "linkedin_url": profile.linkedin_url,
            },
            company={
                "name": profile.company_name,
                "sector": profile.sector,
                "stage": profile.stage,
                "geography": profile.geography,
            },
        )

    # ── stage 1: extract raw text per slide/page ──────────────────────────────
    def _extract_text(self, file_path: str) -> list[str]:
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".pdf":
            return self._from_pdf(file_path)
        if ext == ".pptx":
            return self._from_pptx(file_path)
        if ext == ".docx":
            return self._from_docx(file_path)
        raise ValueError(f"Unsupported deck format: {ext}")

    @staticmethod
    def _from_pdf(path: str) -> list[str]:
        from pypdf import PdfReader

        reader = PdfReader(path)
        return [(page.extract_text() or "").strip() for page in reader.pages]

    @staticmethod
    def _from_pptx(path: str) -> list[str]:
        from pptx import Presentation

        prs = Presentation(path)
        slides: list[str] = []
        for slide in prs.slides:
            parts = [
                shape.text for shape in slide.shapes if shape.has_text_frame and shape.text
            ]
            slides.append("\n".join(parts).strip())
        return slides

    @staticmethod
    def _from_docx(path: str) -> list[str]:
        from docx import Document

        doc = Document(path)
        return ["\n".join(p.text for p in doc.paragraphs if p.text).strip()]

    # ── stage 2: LLM structured extraction (uses the cheap/fast model) ────────
    @staticmethod
    def _cognify(deck_text: str) -> ExtractedProfile:
        prompt = (
            "Extract the founder identity, company details, and every explicit factual "
            "claim from this pitch deck. For each claim, set `location` to the slide it "
            "came from (e.g. 'slide 4').\n\n"
            f"DECK:\n{deck_text[:20000]}"
        )
        return get_llm().structured(
            prompt, ExtractedProfile, system=_EXTRACT_SYSTEM, fast=True
        )
